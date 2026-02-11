"""
GPGT Document Generation Service
=================================

Flask API that handles:
  - PPTX preprocessing (merge Canva's fragmented text runs)
  - Placeholder substitution ({d.xxx} tags)
  - PPTX → PDF conversion via LibreOffice
  - Multi-page PDF merging

Endpoints:
  POST /preprocess     — Upload a PPTX, get preprocessed version + placeholder report
  POST /render         — Upload a PPTX + booking JSON, get a PDF back
  POST /render-batch   — Upload multiple PPTXs + booking JSON, get a merged PDF back
  GET  /health         — Health check
"""

import io
import os
import re
import shutil
import subprocess
import tempfile
import uuid
from pathlib import Path

from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pptx import Presentation
from lxml import etree
from pypdf import PdfReader, PdfWriter

app = Flask(__name__)
CORS(app)

# ─────────────────────────────────────────────────────────────────
# CONFIG
# ─────────────────────────────────────────────────────────────────

SOFFICE_PATH = os.environ.get("SOFFICE_PATH", "/usr/bin/soffice")
MAX_FILE_SIZE = int(os.environ.get("MAX_FILE_SIZE", 50 * 1024 * 1024))  # 50MB
API_KEY = os.environ.get("GPGT_API_KEY", None)  # Optional auth

NSMAP = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}


# ─────────────────────────────────────────────────────────────────
# AUTH MIDDLEWARE
# ─────────────────────────────────────────────────────────────────

@app.before_request
def check_auth():
    """Optional API key authentication."""
    if not API_KEY:
        return  # No auth configured
    if request.endpoint == 'health':
        return  # Skip auth for health check

    provided = request.headers.get("X-API-Key") or request.args.get("api_key")
    if provided != API_KEY:
        return jsonify({"error": "Unauthorized"}), 401


# ─────────────────────────────────────────────────────────────────
# TEXT FRAME EXTRACTION (handles grouped shapes)
# ─────────────────────────────────────────────────────────────────

def get_all_text_frames(shape):
    """
    Recursively extract all text frames from a shape,
    including those inside grouped shapes.
    """
    frames = []

    # Group shape — recurse into children
    if shape.shape_type is not None and shape.shape_type == 6:
        for child_shape in shape.shapes:
            frames.extend(get_all_text_frames(child_shape))
        return frames

    # Also handle groups that don't report shape_type correctly
    if hasattr(shape, 'shapes'):
        try:
            for child_shape in shape.shapes:
                frames.extend(get_all_text_frames(child_shape))
        except Exception:
            pass

    if shape.has_text_frame:
        frames.append(shape.text_frame)

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                frames.append(cell.text_frame)

    return frames


# ─────────────────────────────────────────────────────────────────
# PREPROCESSING — Merge Canva's fragmented text runs
# ─────────────────────────────────────────────────────────────────

def same_formatting(r1, r2):
    """Check if two runs share identical formatting."""
    try:
        c1 = str(r1.font.color.rgb) if r1.font.color and r1.font.color.rgb else None
        c2 = str(r2.font.color.rgb) if r2.font.color and r2.font.color.rgb else None
        return (
            r1.font.name == r2.font.name
            and r1.font.size == r2.font.size
            and r1.font.bold == r2.font.bold
            and r1.font.italic == r2.font.italic
            and r1.font.underline == r2.font.underline
            and c1 == c2
        )
    except Exception:
        return False


def merge_runs_in_paragraph(paragraph):
    """Merge adjacent runs with identical formatting."""
    runs = list(paragraph.runs)
    if len(runs) <= 1:
        return 0
    p_elem = paragraph._p
    r_elements = p_elem.findall('.//a:r', NSMAP)
    if not r_elements or len(r_elements) != len(runs):
        return 0
    groups, group = [], [0]
    for i in range(1, len(runs)):
        if same_formatting(runs[i - 1], runs[i]):
            group.append(i)
        else:
            groups.append(group)
            group = [i]
    groups.append(group)
    removed = 0
    for g in groups:
        if len(g) > 1:
            texts = []
            for idx in g:
                t = r_elements[idx].find('.//a:t', NSMAP)
                texts.append(t.text if t is not None and t.text else '')
            first_t = r_elements[g[0]].find('.//a:t', NSMAP)
            if first_t is not None:
                first_t.text = ''.join(texts)
            for idx in g[1:]:
                p = r_elements[idx].getparent()
                if p is not None:
                    p.remove(r_elements[idx])
                    removed += 1
    return removed


def preprocess_presentation(prs):
    """
    Preprocess a presentation: merge fragmented runs, detect placeholders.
    Returns (merged_count, placeholders, legacy_placeholders).
    """
    total_merged = 0
    placeholders = []
    legacy = []

    for slide in prs.slides:
        for shape in slide.shapes:
            frames = get_all_text_frames(shape)
            for tf in frames:
                for para in tf.paragraphs:
                    total_merged += merge_runs_in_paragraph(para)
                    text = para.text
                    placeholders.extend(re.findall(r'\{d\.\w+\}', text))
                    legacy.extend(re.findall(r'<<\w+>>', text))

    return total_merged, list(set(placeholders)), list(set(legacy))


# ─────────────────────────────────────────────────────────────────
# PLACEHOLDER SUBSTITUTION
# ─────────────────────────────────────────────────────────────────

def substitute_placeholders(prs, booking_data):
    """
    Replace {d.field_name} placeholders with booking data values.
    Two-pass approach: single-run then cross-run for fragmented tags.
    """
    replaced = 0
    unmatched = []
    pattern = re.compile(r'\{d\.(\w+)\}')

    for slide in prs.slides:
        for shape in slide.shapes:
            frames = get_all_text_frames(shape)

            for tf in frames:
                for para in tf.paragraphs:
                    runs = list(para.runs)
                    if not runs:
                        continue

                    # Pass 1: Single-run replacements
                    for run in runs:
                        if not run.text or '{d.' not in run.text:
                            continue

                        def replace_match(match):
                            nonlocal replaced
                            field = match.group(1)
                            value = booking_data.get(field)
                            if value is not None:
                                replaced += 1
                                return str(value)
                            else:
                                unmatched.append(field)
                                return match.group(0)

                        run.text = pattern.sub(replace_match, run.text)

                    # Pass 2: Cross-run replacements
                    runs = list(para.runs)
                    full_text = ''.join(r.text or '' for r in runs)

                    if '{d.' not in full_text:
                        continue

                    for match in pattern.finditer(full_text):
                        field = match.group(1)
                        value = booking_data.get(field)
                        tag = match.group(0)

                        if value is None:
                            if field not in unmatched:
                                unmatched.append(field)
                            continue

                        tag_start = match.start()
                        tag_end = match.end()

                        pos = 0
                        run_ranges = []
                        for ri, run in enumerate(runs):
                            rtext = run.text or ''
                            rlen = len(rtext)
                            run_start = pos
                            run_end = pos + rlen
                            overlap_start = max(tag_start, run_start)
                            overlap_end = min(tag_end, run_end)
                            if overlap_start < overlap_end:
                                local_start = overlap_start - run_start
                                local_end = overlap_end - run_start
                                run_ranges.append((ri, local_start, local_end))
                            pos = run_end

                        if len(run_ranges) <= 1:
                            continue

                        for idx, (ri, local_start, local_end) in enumerate(run_ranges):
                            rtext = runs[ri].text or ''
                            if idx == 0:
                                runs[ri].text = rtext[:local_start] + str(value) + rtext[local_end:]
                            else:
                                runs[ri].text = rtext[:local_start] + rtext[local_end:]

                        replaced += 1
                        full_text = ''.join(r.text or '' for r in runs)
                        if '{d.' not in full_text:
                            break

    return replaced, list(set(unmatched))


# ─────────────────────────────────────────────────────────────────
# PPTX → PDF CONVERSION
# ─────────────────────────────────────────────────────────────────

def pptx_to_pdf(pptx_path, output_dir):
    """Convert PPTX to PDF using LibreOffice headless."""
    env = os.environ.copy()
    env["HOME"] = tempfile.gettempdir()

    cmd = [
        SOFFICE_PATH,
        "--headless",
        "--norestore",
        "--convert-to", "pdf",
        "--outdir", str(output_dir),
        str(pptx_path),
    ]

    result = subprocess.run(
        cmd,
        capture_output=True,
        text=True,
        timeout=120,
        env=env,
    )

    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice conversion failed: {result.stderr[:500]}")

    pdf_name = Path(pptx_path).stem + ".pdf"
    pdf_path = Path(output_dir) / pdf_name

    if not pdf_path.exists():
        raise RuntimeError(f"PDF not generated. stdout: {result.stdout[:300]}")

    return str(pdf_path)


# ─────────────────────────────────────────────────────────────────
# PDF MERGING
# ─────────────────────────────────────────────────────────────────

def merge_pdfs(pdf_paths):
    """Merge multiple PDFs into one, return bytes."""
    writer = PdfWriter()
    for pdf_path in pdf_paths:
        reader = PdfReader(pdf_path)
        for page in reader.pages:
            writer.add_page(page)

    buf = io.BytesIO()
    writer.write(buf)
    buf.seek(0)
    return buf, len(writer.pages)


# ─────────────────────────────────────────────────────────────────
# FULL RENDER PIPELINE
# ─────────────────────────────────────────────────────────────────

def render_single_page(pptx_bytes, booking_data, work_dir):
    """
    Full pipeline for one page:
      1. Load PPTX
      2. Preprocess (merge fragmented runs)
      3. Substitute placeholders
      4. Convert to PDF
    Returns (pdf_path, stats_dict).
    """
    # Save to temp file
    uid = uuid.uuid4().hex[:8]
    pptx_path = os.path.join(work_dir, f"page_{uid}.pptx")
    with open(pptx_path, "wb") as f:
        f.write(pptx_bytes)

    # Load and preprocess
    prs = Presentation(pptx_path)
    merged, placeholders, legacy = preprocess_presentation(prs)

    # Substitute
    replaced, unmatched = substitute_placeholders(prs, booking_data)

    # Save modified PPTX
    modified_path = os.path.join(work_dir, f"modified_{uid}.pptx")
    prs.save(modified_path)

    # Convert to PDF
    pdf_path = pptx_to_pdf(modified_path, work_dir)

    stats = {
        "merged_runs": merged,
        "placeholders_found": placeholders,
        "placeholders_replaced": replaced,
        "unmatched": unmatched,
        "legacy_tags": legacy,
        "slide_count": len(prs.slides),
    }

    return pdf_path, stats


# ─────────────────────────────────────────────────────────────────
# ENDPOINTS
# ─────────────────────────────────────────────────────────────────

@app.route("/health", methods=["GET"])
def health():
    """Health check — also verifies LibreOffice is available."""
    lo_available = os.path.isfile(SOFFICE_PATH) or shutil.which(SOFFICE_PATH) is not None
    return jsonify({
        "status": "ok",
        "libreoffice": lo_available,
        "soffice_path": SOFFICE_PATH,
    })


@app.route("/preprocess", methods=["POST"])
def preprocess_endpoint():
    """
    Upload a raw Canva PPTX, get back:
      - Preprocessed PPTX file (as download)
      - Placeholder report (as JSON)

    Usage:
      POST /preprocess
      Content-Type: multipart/form-data
      Body: template (file)

    Query params:
      ?report_only=true — return JSON report without the file
    """
    if "template" not in request.files:
        return jsonify({"error": "No 'template' file provided"}), 400

    file = request.files["template"]
    pptx_bytes = file.read()

    if len(pptx_bytes) > MAX_FILE_SIZE:
        return jsonify({"error": f"File too large (max {MAX_FILE_SIZE // 1024 // 1024}MB)"}), 413

    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
        merged, placeholders, legacy = preprocess_presentation(prs)

        report = {
            "filename": file.filename,
            "slide_count": len(prs.slides),
            "merged_runs": merged,
            "placeholders": sorted(placeholders),
            "legacy_placeholders": sorted(legacy),
        }

        if request.args.get("report_only") == "true":
            return jsonify(report)

        # Return preprocessed file
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        return send_file(
            buf,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            download_name=f"preprocessed_{file.filename}",
            as_attachment=True,
        )

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/render", methods=["POST"])
def render_endpoint():
    """
    Render a single PPTX page with booking data, return PDF.

    Usage:
      POST /render
      Content-Type: multipart/form-data
      Body:
        template (file) — the PPTX page
        data (text/json) — booking data JSON string

    Returns: PDF file
    """
    if "template" not in request.files:
        return jsonify({"error": "No 'template' file provided"}), 400
    if "data" not in request.form:
        return jsonify({"error": "No 'data' field provided"}), 400

    try:
        pptx_bytes = request.files["template"].read()
        booking_data = __import__("json").loads(request.form["data"])
    except Exception as e:
        return jsonify({"error": f"Invalid input: {e}"}), 400

    try:
        with tempfile.TemporaryDirectory() as work_dir:
            pdf_path, stats = render_single_page(pptx_bytes, booking_data, work_dir)

            return send_file(
                pdf_path,
                mimetype="application/pdf",
                download_name="rendered.pdf",
                as_attachment=True,
            )
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/render-batch", methods=["POST"])
def render_batch_endpoint():
    """
    Render multiple PPTX pages with booking data, return a single merged PDF.
    This is the main endpoint used by the Supabase Edge Function.

    Usage:
      POST /render-batch
      Content-Type: multipart/form-data
      Body:
        pages (files) — multiple PPTX files, uploaded in order
        data (text/json) — booking data JSON string

    Returns: Single merged PDF
    """
    files = request.files.getlist("pages")
    if not files:
        return jsonify({"error": "No 'pages' files provided"}), 400
    if "data" not in request.form:
        return jsonify({"error": "No 'data' field provided"}), 400

    try:
        booking_data = __import__("json").loads(request.form["data"])
    except Exception as e:
        return jsonify({"error": f"Invalid booking data JSON: {e}"}), 400

    try:
        with tempfile.TemporaryDirectory() as work_dir:
            pdf_paths = []
            all_stats = []

            for i, file in enumerate(files):
                pptx_bytes = file.read()
                pdf_path, stats = render_single_page(pptx_bytes, booking_data, work_dir)
                pdf_paths.append(pdf_path)
                all_stats.append({
                    "page": i + 1,
                    "filename": file.filename,
                    **stats,
                })

            # Merge all PDFs
            merged_buf, total_pages = merge_pdfs(pdf_paths)

            return send_file(
                merged_buf,
                mimetype="application/pdf",
                download_name="final_document.pdf",
                as_attachment=True,
            )

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/render-batch-from-urls", methods=["POST"])
def render_batch_from_urls_endpoint():
    """
    Render multiple PPTX pages (fetched from URLs) with booking data.
    This is the endpoint called by the Supabase Edge Function.

    The Edge Function sends URLs to PPTX files in Supabase Storage
    rather than uploading files directly.

    Usage:
      POST /render-batch-from-urls
      Content-Type: application/json
      Body:
        {
          "pages": [
            {"url": "https://xxx.supabase.co/storage/v1/object/public/page-templates/cover.pptx"},
            {"url": "https://xxx.supabase.co/storage/v1/object/public/page-templates/flights_standard.pptx"},
            ...
          ],
          "booking_data": {
            "lead_name": "James Mitchell",
            ...
          }
        }

    Returns: PDF file
    """
    import requests as http_requests

    try:
        body = request.get_json()
        pages = body.get("pages", [])
        booking_data = body.get("booking_data", {})
    except Exception as e:
        return jsonify({"error": f"Invalid JSON body: {e}"}), 400

    if not pages:
        return jsonify({"error": "No pages provided"}), 400

    try:
        with tempfile.TemporaryDirectory() as work_dir:
            pdf_paths = []

            for i, page in enumerate(pages):
                url = page.get("url")
                if not url:
                    continue

                # Download PPTX from URL
                resp = http_requests.get(url, timeout=30)
                resp.raise_for_status()

                pdf_path, stats = render_single_page(resp.content, booking_data, work_dir)
                pdf_paths.append(pdf_path)

            if not pdf_paths:
                return jsonify({"error": "No pages were rendered successfully"}), 500

            # Merge
            merged_buf, total_pages = merge_pdfs(pdf_paths)

            return send_file(
                merged_buf,
                mimetype="application/pdf",
                download_name="final_document.pdf",
                as_attachment=True,
            )

    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ─────────────────────────────────────────────────────────────────
# RUN
# ─────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8080))
    debug = os.environ.get("FLASK_DEBUG", "false").lower() == "true"
    app.run(host="0.0.0.0", port=port, debug=debug)
